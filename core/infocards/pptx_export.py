"""簡報 → PPTX 匯出（從 infoCard utils/pptExporter.ts + data/slideMasters.ts 收編，Phase C）。

slideMasters.ts 是「PPTX 匯出器與 React renderer 共用」的座標真相（10" × 5.625" 畫布），
本模組直接依該座標用 python-pptx 渲染，確保輸出與前端預覽一致。

PoC 範圍：涵蓋全 19 版型的「內容承載元素」（accent bar + 標題 + 主體：bullets / statValue /
雙欄 / 圖片 / chart_focus 原生圖表 / 教學版型 content+steps+answer / code）。裝飾性元素
（圓圈、引號大字、時間軸線）暫略——不影響可讀的可下載 .pptx。重依賴（pptx）lazy import。
"""
from __future__ import annotations

import base64
import io

# 畫布（對齊 slideMasters.ts SLIDE）。
SLIDE_W = 10.0
SLIDE_H = 5.625

# 共用幾何（對齊 slideMasters.ts G）。
_MARGIN = 0.45
_HEADER_Y = 0.28
_HEADER_H = 0.72
_CONTENT_Y = 1.15
_CONTENT_H = 3.85
_CONTENT_W = 9.2
_SIDEBAR_W = 0.10

_DEFAULT_ACCENT = "1e40af"


def _hex(color: str | None, fallback: str = _DEFAULT_ACCENT) -> str:
    """'#1e40af' / '1e40af' → 6 碼 hex（給 RGBColor.from_string）；非法退 fallback。"""
    if not color:
        return fallback
    c = color.lstrip("#").strip()
    if len(c) == 3:  # #abc → aabbcc
        c = "".join(ch * 2 for ch in c)
    if len(c) != 6:
        return fallback
    try:
        int(c, 16)
        return c.lower()
    except ValueError:
        return fallback


def _decode_data_url(data_url: str) -> bytes | None:
    """data:image/png;base64,xxxx → bytes；非 base64 data URL 回 None。"""
    if not data_url or "base64," not in data_url:
        return None
    try:
        return base64.b64decode(data_url.split("base64,", 1)[1])
    except Exception:
        return None


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False, italic=False,
              align="left", valign="top", color="111111"):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    lines = text if isinstance(text, list) else str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = str(line)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_rect(slide, x, y, w, h, color):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(color)
    shp.line.fill.background()
    return shp


def _add_image(slide, x, y, w, h, data_url) -> bool:
    from pptx.util import Inches

    raw = _decode_data_url(data_url)
    if not raw:
        return False
    try:
        slide.shapes.add_picture(io.BytesIO(raw), Inches(x), Inches(y), Inches(w), Inches(h))
        return True
    except Exception:
        return False


def _bullets(slide_obj) -> list[str]:
    bp = slide_obj.get("bulletPoints") or []
    return [str(b) for b in bp if b]


def _add_chart(slide, x, y, w, h, chart_data) -> bool:
    """chart_focus 原生圖表（bar→COLUMN_CLUSTERED / pie→PIE）。資料不足回 False。"""
    if not isinstance(chart_data, dict):
        return False
    labels = chart_data.get("labels") or []
    values = chart_data.get("values") or []
    if len(labels) < 2 or len(labels) != len(values):
        return False
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches

        cd = CategoryChartData()
        cd.categories = [str(l) for l in labels]
        cd.add_series("數值", tuple(float(v) for v in values))
        ctype = XL_CHART_TYPE.PIE if chart_data.get("type") == "pie" else XL_CHART_TYPE.COLUMN_CLUSTERED
        slide.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), cd)
        return True
    except Exception:
        return False


def _render_slide(prs, s: dict, accent: str) -> None:
    from pptx.dml.color import RGBColor

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    layout = s.get("layout", "bullet_list")
    title = s.get("title", "")
    content = s.get("content", "") or ""
    has_image = bool(s.get("imageUrl"))

    # 內容版型左側 accent bar（對齊 masters accentBar；結構性版型如封面不畫）。
    _STRUCTURAL = {"title_cover", "section_header", "conclusion", "diagram_image", "full_image"}
    if layout not in _STRUCTURAL:
        _add_rect(slide, 0, 0, _SIDEBAR_W, 5.3, accent)

    if layout == "title_cover":
        _add_text(slide, 0.8, 1.0, 8.4, 2.2, title, size=36, bold=True, align="center", valign="middle", color=accent)
        if s.get("subtitle") or content:
            _add_text(slide, 1.5, 3.5, 7.0, 1.4, s.get("subtitle") or content, size=18, align="center")
        return
    if layout == "section_header":
        _add_text(slide, 0.8, 1.6, 8.4, 1.6, title, size=42, bold=True, align="center", color=accent)
        if content:
            _add_text(slide, 1.5, 3.4, 7.0, 0.9, content, size=20, align="center")
        return
    if layout == "conclusion":
        _add_text(slide, 0.5, 1.1, 9.0, 0.9, title, size=34, bold=True, align="center", color=accent)
        body = content or "\n".join(_bullets(s))
        if body:
            _add_text(slide, 1.0, 2.35, 8.0, 2.55, body, size=18, align="center")
        return

    # 一般內容版型：頂部標題 header。
    _add_text(slide, _MARGIN, _HEADER_Y, _CONTENT_W, _HEADER_H, title, size=27, bold=True, color=accent)

    if layout == "big_number":
        if s.get("statValue"):
            _add_text(slide, 2.2, 1.35, 5.6, 2.4, s["statValue"], size=58, bold=True, align="center", valign="middle", color=accent)
        if content:
            _add_text(slide, 1.0, 3.95, 8.0, 0.95, content, size=17, align="center")
        return
    if layout == "two_column":
        left, right = s.get("columnLeft") or [], s.get("columnRight") or []
        labels = (content or "|").split("|")
        ltitle = labels[0] if labels else ""
        rtitle = labels[1] if len(labels) > 1 else ""
        _add_text(slide, _MARGIN, 1.2, 4.35, 3.8,
                  ([ltitle] if ltitle else []) + [f"• {x}" for x in left], size=15)
        _add_text(slide, 5.35, 1.2, 4.3, 3.8,
                  ([rtitle] if rtitle else []) + [f"• {x}" for x in right], size=15)
        return
    if layout in ("text_and_image", "diagram_image", "full_image"):
        placed = has_image and _add_image(
            slide, *({"text_and_image": (5.7, 1.1, 3.9, 3.9),
                      "diagram_image": (0.3, 0.8, 9.4, 4.2),
                      "full_image": (0.0, 0.0, SLIDE_W, SLIDE_H)}[layout]), s.get("imageUrl"))
        body = content or "\n".join(_bullets(s))
        if layout == "text_and_image":
            _add_text(slide, _MARGIN, _CONTENT_Y, 5.0 if placed else _CONTENT_W, _CONTENT_H, body, size=16)
        elif not placed and body:
            _add_text(slide, 1.0, 2.0, 8.0, 2.0, body, size=18, align="center")
        return
    if layout == "chart_focus":
        if not _add_chart(slide, 0.4, _CONTENT_Y, 9.2, 3.7, s.get("chartData")):
            _add_text(slide, 0.4, _CONTENT_Y, 9.2, 3.7, content or "\n".join(_bullets(s)), size=18)
        return
    if layout == "code_block":
        _add_rect(slide, 0.25, 1.05, 6.2, 4.0, "1e293b")
        _add_text(slide, 0.35, 1.15, 6.0, 3.85, content, size=11, color="e2e8f0")
        if _bullets(s):
            _add_text(slide, 6.6, 1.15, 3.15, 3.85, [f"• {x}" for x in _bullets(s)], size=13)
        return
    if layout in ("worked_example", "exercise"):
        _add_text(slide, 0.5, 1.15, 9.0, 1.2, content, size=15)
        steps = _bullets(s)
        if steps:
            _add_text(slide, 0.5, 2.4, 9.0, 2.3, [f"{i+1}. {x}" for i, x in enumerate(steps)], size=14)
        if s.get("statValue"):
            _add_text(slide, 5.5, 4.1, 4.2, 0.85, f"Ans: {s['statValue']}", size=18, bold=True, align="center", valign="middle", color=accent)
        return
    if layout == "quote":
        _add_text(slide, 0.85, 1.2, 8.0, 2.6, content or title, size=22, italic=True, valign="middle")
        if title and content:
            _add_text(slide, 4.5, 4.08, 5.1, 0.6, f"— {title}", size=15, bold=True, align="right")
        return
    if layout == "swot_analysis":
        q = s.get("quadrants") or {}
        cells = [("優勢 Strengths", q.get("strengths"), 0.45, 1.15), ("劣勢 Weaknesses", q.get("weaknesses"), 5.05, 1.15),
                 ("機會 Opportunities", q.get("opportunities"), 0.45, 3.25), ("威脅 Threats", q.get("threats"), 5.05, 3.25)]
        for label, items, cx, cy in cells:
            _add_text(slide, cx, cy, 4.5, 1.95, [label] + [f"• {x}" for x in (items or [])], size=13)
        return
    if layout == "pyramid_diagram":
        layers = s.get("pyramidLayers") or []
        _add_text(slide, 1.5, 1.15, 7.0, 4.0, [f"▸ {x}" for x in layers] or [content], size=16, align="center", valign="middle")
        return
    if layout == "comparison_table":
        cmp = s.get("comparisonData") or {}
        _render_table(slide, cmp.get("headers") or [], cmp.get("rows") or [], accent)
        return

    # 預設（bullet_list / process_steps / timeline / 其他）：標題 + 條列。
    body = [f"• {x}" for x in _bullets(s)] or [content]
    _add_text(slide, 0.6, _CONTENT_Y, 9.0, _CONTENT_H, body, size=18)


def _render_table(slide, headers, rows, accent) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    if not headers and not rows:
        return
    ncol = max(len(headers), max((len(r) for r in rows), default=0)) or 1
    nrow = (1 if headers else 0) + len(rows)
    if nrow == 0:
        return
    tbl = slide.shapes.add_table(nrow, ncol, Inches(0.45), Inches(1.15), Inches(9.1), Inches(3.8)).table
    r0 = 0
    if headers:
        for c in range(ncol):
            cell = tbl.cell(0, c)
            cell.text = str(headers[c]) if c < len(headers) else ""
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(accent)
        r0 = 1
    for i, row in enumerate(rows):
        for c in range(ncol):
            tbl.cell(r0 + i, c).text = str(row[c]) if c < len(row) else ""


def build_pptx(data: dict) -> bytes:
    """PresentationData(dict) → .pptx bytes。空 slides 也產出僅封面的最小檔。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    accent = _hex(data.get("themeColor"))

    slides = data.get("slides") or []
    if not slides:
        slides = [{"layout": "title_cover", "title": data.get("mainTitle", "簡報"),
                   "subtitle": data.get("subtitle", "")}]
    for s in slides:
        try:
            _render_slide(prs, s, accent)
        except Exception:
            # 單頁渲染失敗不該毀掉整份匯出：退成標題頁。
            _render_slide(prs, {"layout": "section_header", "title": s.get("title", "（此頁無法渲染）")}, accent)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
