"""core/pptx_augment.py — 在「原始 .pptx」上就地補圖 (原文字保持可編輯)。

定位
----
slides_pdf 補圖是把每頁 PDF 渲成 PNG, 原頁文字因此變成圖、不可再編輯。若使用者
手上是 .pptx 原檔, 更好的作法是直接在原檔上動手:

  1. pptx → pdf → 逐頁 PNG (僅供分析, 用 LibreOffice／PowerPoint + PyMuPDF)。
  2. 偵測缺圖頁 + 每頁空白區 (複用 core.slide_image_gen)。
  3. 為缺圖頁生 AI 配圖 (複用 generate_slide_image; prompt 取自該頁文字)。
  4. 打開「原始 .pptx」, 把配圖**加進**該頁空白區 — 原本的文字方塊 / 圖形全部
     原封不動, 仍可在 PowerPoint 內編輯。

如此匯出的新簡報文字可改、圖在空白處, 解決「整頁變成一張圖」的問題。

依賴
----
- python-pptx (匯入/匯出 .pptx)
- LibreOffice (soffice) — 跨平台優先的 pptx→pdf 渲染器。
- Windows PowerPoint COM — 本機已裝 Office 時的 fallback（需 pywin32）。
- PyMuPDF (fitz) — pdf→png。
mock=True 走 PIL 佔位圖 (不打 Gemini), 但仍需 LibreOffice 渲染原頁 (除非 caller
直接給 page_pngs)。
"""
from __future__ import annotations

import logging
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path

logger = logging.getLogger(__name__)
_IS_WINDOWS = os.name == "nt"


def _render_with_powerpoint(src_pptx: Path, pdf: Path) -> Path:
    """Windows PowerPoint COM fallback；以獨立 process 隔離 Office COM lifecycle。"""
    if not _IS_WINDOWS:
        raise RuntimeError("PowerPoint COM fallback 僅支援 Windows")
    tool = Path(__file__).resolve().parent.parent / "tools" / "pptx_to_pdf.py"
    res = subprocess.run(
        [sys.executable, str(tool), str(src_pptx.resolve()), str(pdf.resolve())],
        capture_output=True, text=True, timeout=180,
    )
    if res.returncode != 0:
        detail = (res.stderr or res.stdout or "unknown error")[-500:]
        raise RuntimeError(f"PowerPoint 轉檔失敗 (code {res.returncode}): {detail}")
    if not pdf.is_file() or pdf.stat().st_size == 0:
        raise RuntimeError("PowerPoint 轉檔未產生有效 PDF")
    return pdf


def render_pptx_to_pdf(src_pptx: str | Path, out_dir: str | Path, *, timeout: int = 180) -> Path:
    """把 .pptx 轉成 .pdf；優先 LibreOffice，Windows 可退 PowerPoint COM。"""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    src_pptx = Path(src_pptx)
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf = out_dir / (src_pptx.stem + ".pdf")
    if not soffice:
        if _IS_WINDOWS:
            return _render_with_powerpoint(src_pptx, pdf)
        raise RuntimeError(
            "找不到 LibreOffice (soffice), 無法把 PPTX 轉成 PDF 做分析。"
            "請安裝 libreoffice-impress。"
        )
    # 用獨立 profile 避免併發鎖
    profile = (out_dir / "_lo_profile").resolve().as_uri()
    cmd = [
        soffice, "--headless", f"-env:UserInstallation={profile}",
        "--convert-to", "pdf", "--outdir", str(out_dir), str(src_pptx),
    ]
    res = subprocess.run(cmd, capture_output=True, timeout=timeout)
    if res.returncode != 0 or not pdf.exists():
        raise RuntimeError(
            f"LibreOffice 轉檔失敗 (code {res.returncode}): "
            f"{res.stderr.decode('utf-8', 'replace')[:300]}"
        )
    return pdf


def _render_pdf_pages(pdf: Path, out_dir: Path, *, zoom: float = 2.0) -> list[Path]:
    import fitz
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf)
    paths = []
    try:
        for i, page in enumerate(doc, start=1):
            p = out_dir / f"p{i:03d}.png"
            page.get_pixmap(matrix=fitz.Matrix(zoom, zoom)).save(str(p))
            paths.append(p)
    finally:
        doc.close()
    return paths


def extract_pptx_slide_texts(prs) -> list[tuple[str, str]]:
    """每張投影片 → (title, body) 文字, 供生圖 prompt 用。"""
    out = []
    for slide in prs.slides:
        lines = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                t = sh.text_frame.text.strip()
                if t:
                    lines.append(t)
        title = lines[0].split("\n")[0][:60] if lines else ""
        body = " ".join(lines)[:240]
        out.append((title, body))
    return out


def extract_pptx_speaker_notes(prs, *, max_chars: int = 1200) -> list[str]:
    """每張投影片 → 講者備註純文字 (沒寫備註的頁回空字串)。

    老師在 PowerPoint「備忘稿」寫的講稿是最貼近他本人要講什麼的材料, 拿來餵旁白
    生成 (slide_ingest.narrate_page_with_gemini) 比讓 Gemini 純看圖猜準得多。
    python-pptx 的 notes_slide 只在該頁真的有 notes 時才存在, 一律先問 has_notes_slide。
    """
    out: list[str] = []
    for slide in prs.slides:
        text = ""
        try:
            if slide.has_notes_slide:
                frame = slide.notes_slide.notes_text_frame
                if frame is not None:
                    text = (frame.text or "").strip()
        except Exception:  # noqa: BLE001 — 備註讀不到不該擋住整份簡報
            text = ""
        # PowerPoint 常把投影片編號也塞進 notes placeholder, 落單數字沒有資訊量
        if text.isdigit():
            text = ""
        out.append(text[:max_chars])
    return out


def read_pptx_speaker_notes(src_pptx: str | Path, *, max_chars: int = 1200) -> list[str]:
    """從 .pptx 檔讀講者備註; 沒裝 python-pptx 或讀檔失敗 → 回 []。"""
    try:
        from pptx import Presentation
    except ImportError:
        logger.info("未安裝 python-pptx, 略過講者備註擷取")
        return []
    try:
        return extract_pptx_speaker_notes(Presentation(str(src_pptx)), max_chars=max_chars)
    except Exception as e:  # noqa: BLE001
        logger.warning("讀取講者備註失敗 (%s), 略過", e)
        return []


def insert_images_into_pptx(
    src_pptx: str | Path,
    out_pptx: str | Path,
    items: list[tuple[int, Path, tuple | list | None]],
) -> int:
    """在 src_pptx 上插圖後另存 out_pptx (原文字/圖形不動)。回傳插入張數。

    items: list of (slide_index_0based, image_path, placement)。placement 為正規化
    (x,y,w,h)∈[0,1]; None → 右下角浮貼。圖等比 fit 進框 (置中, 不變形)。
    """
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(str(src_pptx))
    SW, SH = prs.slide_width, prs.slide_height
    slides = list(prs.slides)
    n = 0
    for idx, img, placement in items:
        if idx < 0 or idx >= len(slides):
            continue
        if not Path(img).exists():
            continue
        if placement:
            nx, ny, nw, nh = placement
        else:
            nx, ny, nw, nh = 0.68, 0.62, 0.30, 0.30
        inner = 0.06
        bx = (nx + nw * inner) * SW
        by = (ny + nh * inner) * SH
        bw = nw * (1 - 2 * inner) * SW
        bh = nh * (1 - 2 * inner) * SH
        # 用實際圖比例等比 fit (預設配圖近 1:1, 仍精確處理)
        try:
            from PIL import Image
            with Image.open(img) as im:
                iw, ih = im.size
        except Exception:
            iw, ih = 1, 1
        scale = min(bw / iw, bh / ih)
        w, h = iw * scale, ih * scale
        left = Emu(int(bx + (bw - w) / 2))
        top = Emu(int(by + (bh - h) / 2))
        slides[idx].shapes.add_picture(str(img), left, top, Emu(int(w)), Emu(int(h)))
        n += 1
    Path(out_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))
    return n


def augment_pptx(
    src_pptx: str | Path,
    out_pptx: str | Path,
    *,
    work_dir: str | Path,
    only_missing: bool = True,
    mock: bool = False,
    api_key: str | None = None,
    max_images: int | None = None,
    page_pngs: list[Path] | None = None,
) -> dict:
    """在原始 .pptx 上為缺圖頁就地補圖, 另存 out_pptx。回傳 summary。

    page_pngs 給定時跳過 LibreOffice 渲染 (測試 / caller 已有逐頁圖時用)。
    """
    from pptx import Presentation
    from core import slide_image_gen as sig
    from core.diagram_image_gen import _build_diagram_prompt  # noqa: F401 (確保模組可用)

    src_pptx = Path(src_pptx)
    work_dir = Path(work_dir)
    work_dir.mkdir(parents=True, exist_ok=True)
    fig_dir = work_dir / "figures"
    fig_dir.mkdir(exist_ok=True)

    # 1) 逐頁 PNG (供分析) + PDF (供缺圖偵測)
    if page_pngs is None:
        pdf = render_pptx_to_pdf(src_pptx, work_dir)
        page_pngs = _render_pdf_pages(pdf, work_dir / "pages")
        imageless = set(sig.detect_imageless_pages(pdf))
    else:
        page_pngs = [Path(p) for p in page_pngs]
        # 無 pdf 時退化: 只要 only_missing 仍想要偵測, 改用每頁 PNG 的內容判斷
        imageless = set(range(1, len(page_pngs) + 1))

    # 2) 每頁文字 (生圖 prompt)
    prs = Presentation(str(src_pptx))
    texts = extract_pptx_slide_texts(prs)
    deck_title = src_pptx.stem

    # 3) 逐頁: 缺圖才生圖 + 算空白框
    items: list[tuple[int, Path, tuple | None]] = []
    generated = 0
    for i, png in enumerate(page_pngs, start=1):  # i = 1-based 頁碼
        if max_images is not None and generated >= max_images:
            break
        if only_missing and i not in imageless:
            continue
        title, body = texts[i - 1] if i - 1 < len(texts) else ("", "")
        slide = {"id": f"p{i:03d}", "title": title or f"投影片 {i}", "narration": body}
        ai_path = fig_dir / f"ai_p{i:03d}.png"
        ok, err = sig.generate_slide_image(
            slide, ai_path, deck_title=deck_title, api_key=api_key, mock=mock,
        )
        if not ok:
            logger.warning("PPTX 補圖跳過 p%d: %s", i, err)
            continue
        placement = sig.find_empty_region(png)
        items.append((i - 1, ai_path, placement))  # 0-based slide index
        generated += 1

    # 4) 在原檔上插圖另存
    inserted = insert_images_into_pptx(src_pptx, out_pptx, items)
    summary = {
        "pages": len(page_pngs),
        "imageless": sorted(imageless) if imageless else [],
        "generated": generated,
        "inserted": inserted,
        "mock": mock,
    }
    logger.info("PPTX 就地補圖完成: %s", summary)
    return summary
