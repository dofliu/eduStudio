"""core/slide_pptx.py — 把 slides deck 匯出成 .pptx (含 AI 補圖)。

定位
----
slide_image_gen 為缺圖簡報逐頁補圖後, deck 的 slide 有:
  - bg_image:        合成後的新頁 (原頁+配圖)
  - source_bg_image: 原頁 (補圖前)
  - ai_image:        AI 配圖
  - narration:       逐頁旁白

這個模組把整份 deck 匯出成可編輯 .pptx (16:9):
  - 補過圖的頁 → 依 layout 放「原頁 + AI 配圖」兩張**獨立**圖片 (PowerPoint 內
    可各自移動 / 縮放 / 換圖), 而非烤死的合成圖。
  - 沒補圖的頁 → 原頁鋪滿。
  - 每頁 narration → 講者備忘稿 (speaker notes)。

跟 core/infocards/pptx_export.py 的差別: 那個是 infoCard 卡片 (data-URL 圖、19 種
版型); 這個專餵 slides deck (檔案路徑圖、逐頁一張投影片)。

python-pptx 是 optional dep (lazy import); 沒裝 → raise RuntimeError 提示安裝。
"""
from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# 16:9 投影片 (吋)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
PAD_IN = 0.25


def _img_size(path: Path) -> tuple[int, int] | None:
    try:
        from PIL import Image
        with Image.open(path) as im:
            return im.size
    except Exception as e:  # noqa: BLE001
        logger.warning("讀圖尺寸失敗 %s: %s", path, e)
        return None


def _fit_box(path: Path, bx: float, by: float, bw: float, bh: float):
    """把 path 的圖等比 letterbox-fit 進 (bx,by,bw,bh) 吋方框, 回 (left,top,w,h) 吋。

    讀不到尺寸時退化成填滿方框 (含 pad)。
    """
    size = _img_size(path)
    avail_w, avail_h = bw - 2 * PAD_IN, bh - 2 * PAD_IN
    if not size:
        return (bx + PAD_IN, by + PAD_IN, avail_w, avail_h)
    iw, ih = size
    scale = min(avail_w / iw, avail_h / ih)
    w, h = iw * scale, ih * scale
    return (bx + (bw - w) / 2, by + (bh - h) / 2, w, h)


def _add_picture(slide, path: Path, box) -> bool:
    from pptx.util import Inches
    left, top, w, h = box
    try:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))
        return True
    except Exception as e:  # noqa: BLE001
        logger.warning("PPTX 加圖失敗 %s: %s", path, e)
        return False


def _place_slide_images(
    slide, *, original: Path | None, ai: Path | None, layout: str,
    placement: tuple | list | None = None,
) -> None:
    """依 layout 把原頁 / 配圖放進一張 PPTX slide。

    layout="auto" 且有 placement (正規化 x,y,w,h) → 原頁鋪滿, 配圖就地置入該空白框
    (兩者皆為獨立可編輯圖片)。
    """
    W, H, half = SLIDE_W_IN, SLIDE_H_IN, SLIDE_W_IN / 2
    has_orig = bool(original and original.exists())
    has_ai = bool(ai and ai.exists())

    if layout == "auto" and has_orig and has_ai:
        _add_picture(slide, original, _fit_box(original, 0, 0, W, H))
        if placement:
            nx, ny, nw, nh = placement
            inner = 0.06
            bx = (nx + nw * inner) * W
            by = (ny + nh * inner) * H
            bw = nw * (1 - 2 * inner) * W
            bh = nh * (1 - 2 * inner) * H
        else:  # 無 placement → 右下角浮貼
            bw, bh = W * 0.30, H * 0.30
            bx, by = W - bw - W * 0.02, H - bh - H * 0.02
        _add_picture(slide, ai, _fit_box(ai, bx, by, bw, bh))
        return

    if has_ai and not has_orig or layout == "image_only":
        if has_ai:
            _add_picture(slide, ai, _fit_box(ai, 0, 0, W, H))
        elif has_orig:
            _add_picture(slide, original, _fit_box(original, 0, 0, W, H))
        return

    if not has_ai and has_orig:
        _add_picture(slide, original, _fit_box(original, 0, 0, W, H))
        return

    # 兩張都有 → 依 layout
    if layout == "image_left":
        _add_picture(slide, ai, _fit_box(ai, 0, 0, half, H))
        _add_picture(slide, original, _fit_box(original, half, 0, half, H))
    elif layout == "overlay":
        _add_picture(slide, original, _fit_box(original, 0, 0, W, H))
        iw, ih = W * 0.38, H * 0.38
        _add_picture(slide, ai, _fit_box(ai, W - iw - PAD_IN, H - ih - PAD_IN, iw, ih))
    else:  # side_by_side (預設)
        _add_picture(slide, original, _fit_box(original, 0, 0, half, H))
        _add_picture(slide, ai, _fit_box(ai, half, 0, half, H))


def _set_notes(slide, text: str) -> None:
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except Exception as e:  # noqa: BLE001
        logger.warning("PPTX 寫備忘稿失敗: %s", e)


def deck_to_pptx(deck: dict, out_path: str | Path, *, asset_base: str | Path) -> Path:
    """把 slides deck 匯出成 .pptx 寫到 out_path, 回傳 out_path。

    Args:
        deck: slide_ingest / augment 後的 deck (sections/slides)。
        out_path: 輸出 .pptx 路徑。
        asset_base: 解析 slide 圖片相對路徑的基底 (通常 PROJECT_ROOT)。
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as e:
        raise RuntimeError(
            "python-pptx 未安裝, 無法匯出 PPTX。請執行: pip install python-pptx"
        ) from e

    base = Path(asset_base)
    layout = (deck.get("image_augmentation") or {}).get("layout", "side_by_side")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]  # 全空白版面

    n = 0
    for section in deck.get("sections", []):
        for slide in section.get("slides", []):
            s = prs.slides.add_slide(blank)

            def _resolve(rel):
                return (base / rel) if rel else None

            if slide.get("image_generated"):
                original = _resolve(slide.get("source_bg_image"))
                ai = _resolve(slide.get("ai_image"))
                _place_slide_images(
                    s, original=original, ai=ai, layout=layout,
                    placement=slide.get("ai_placement"),
                )
            else:
                bg = _resolve(slide.get("bg_image"))
                if bg and bg.exists():
                    _add_picture(s, bg, _fit_box(bg, 0, 0, SLIDE_W_IN, SLIDE_H_IN))

            _set_notes(s, slide.get("narration") or "")
            n += 1

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    logger.info("PPTX 匯出完成: %d 頁 → %s", n, out_path)
    return out_path
