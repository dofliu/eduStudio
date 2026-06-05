"""core/infocards/pptx_export + /api/export/pptx 測試（Phase C PPTX 匯出 PoC）。

不打 API；用 python-pptx 回讀產出的 .pptx 驗結構。需 pptx 安裝。
"""
from __future__ import annotations

import io

import pytest

pytest.importorskip("pptx", reason="需要 python-pptx")

import core.infocards.pptx_export as px

_DECK = {
    "mainTitle": "控制系統導論",
    "subtitle": "從傳遞函數到穩定性",
    "themeColor": "#1e3a5f",
    "style": "navy",
    "slides": [
        {"id": "s1", "layout": "title_cover", "title": "控制系統導論", "subtitle": "第 1 章"},
        {"id": "s2", "layout": "bullet_list", "title": "大綱",
         "bulletPoints": ["傳遞函數", "穩定性判據", "頻率響應"]},
        {"id": "s3", "layout": "big_number", "title": "關鍵指標", "statValue": "98.6%", "content": "系統穩定度"},
        {"id": "s4", "layout": "two_column", "title": "比較", "content": "優點|缺點",
         "columnLeft": ["快速", "精確"], "columnRight": ["成本高"]},
        {"id": "s5", "layout": "chart_focus", "title": "極點分布",
         "chartData": {"labels": ["P1", "P2", "P3"], "values": [1.0, 2.0, 3.0], "type": "bar"}},
        {"id": "s6", "layout": "worked_example", "title": "解題", "content": "求特徵根",
         "bulletPoints": ["代入 s", "化簡", "求解"], "statValue": "s=-2"},
        {"id": "s7", "layout": "conclusion", "title": "總結", "content": "掌握三大主題"},
    ],
}


class TestHexHelper:
    def test_normal(self):
        assert px._hex("#1e3a5f") == "1e3a5f"
        assert px._hex("1E3A5F") == "1e3a5f"

    def test_shorthand_and_bad(self):
        assert px._hex("#abc") == "aabbcc"
        assert px._hex("zzz") == px._DEFAULT_ACCENT
        assert px._hex(None) == px._DEFAULT_ACCENT


class TestBuildPptx:
    def test_produces_valid_pptx_with_all_slides(self):
        from pptx import Presentation

        blob = px.build_pptx(_DECK)
        assert isinstance(blob, bytes) and len(blob) > 1000
        prs = Presentation(io.BytesIO(blob))
        assert len(prs.slides) == 7  # 每張投影片一頁

    def test_title_text_present(self):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(px.build_pptx(_DECK)))
        all_text = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if shp.has_text_frame:
                    all_text.append(shp.text_frame.text)
        joined = "\n".join(all_text)
        assert "控制系統導論" in joined
        assert "傳遞函數" in joined  # bullet
        assert "98.6%" in joined     # big_number statValue

    def test_chart_focus_creates_chart(self):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(px.build_pptx(_DECK)))
        # 第 5 張 chart_focus 應含 graphic frame chart
        has_chart = any(
            any(getattr(shp, "has_chart", False) for shp in slide.shapes)
            for slide in prs.slides
        )
        assert has_chart

    def test_empty_slides_yields_cover(self):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(px.build_pptx({"mainTitle": "空簡報", "slides": []})))
        assert len(prs.slides) == 1

    def test_dimensions_16_9(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(io.BytesIO(px.build_pptx(_DECK)))
        assert prs.slide_width == Inches(10)
        assert prs.slide_height == Inches(5.625)


class TestExportRoute:
    @pytest.fixture
    def client(self, tmp_path, monkeypatch):
        pytest.importorskip("fastapi.testclient")
        pytest.importorskip("multipart")
        from fastapi.testclient import TestClient

        import server.routes.infocards as ic
        from core.infocards.share_store import ShareStore
        from server.main import create_app

        monkeypatch.setattr(ic, "get_share_store", lambda: ShareStore(db_path=str(tmp_path / "s.db")))
        app = create_app()
        with TestClient(app) as c:
            yield c

    def test_export_returns_pptx(self, client):
        r = client.post("/api/export/pptx", json={"data": _DECK, "filename": "控制系統"})
        assert r.status_code == 200
        assert r.headers["content-type"] == \
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        assert "attachment" in r.headers["content-disposition"]
        assert r.content[:2] == b"PK"  # zip/pptx 魔術位元組
