"""core.photo_deck — 照片 → vision 分析 → deck 測試 (全 mock, 不打 Gemini)。需要 Pillow。"""
from __future__ import annotations

import pytest

pytest.importorskip("PIL", reason="需要 Pillow")

from core import photo_deck as pd


def _photo(path, size=(1600, 1200), color=(120, 60, 200)):
    from PIL import Image
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, color).save(path, format="JPEG")
    return path


class TestDownscale:
    def test_downscale_caps_long_edge(self, tmp_path):
        from PIL import Image
        p = _photo(tmp_path / "big.jpg", size=(4000, 3000))
        raw = pd._downscale_to_bytes(p, max_dim=1024)
        assert raw and raw[:2] == b"\xff\xd8"  # JPEG SOI
        import io
        assert max(Image.open(io.BytesIO(raw)).size) == 1024

    def test_unreadable_returns_none(self, tmp_path):
        bad = tmp_path / "nope.jpg"; bad.write_bytes(b"not an image")
        assert pd._downscale_to_bytes(bad) is None


class TestAnalyzeMock:
    def test_all_kept_with_captions(self, tmp_path):
        paths = [_photo(tmp_path / f"p{i}.jpg") for i in range(3)]
        out = pd.analyze_photos(paths, deck_title_hint="宜蘭之旅", mock=True)
        assert out["title"] == "宜蘭之旅"
        assert len(out["photos"]) == 3
        assert all(p["keep"] and p["caption"] for p in out["photos"])

    def test_max_select_trims(self, tmp_path):
        paths = [_photo(tmp_path / f"p{i}.jpg") for i in range(5)]
        out = pd.analyze_photos(paths, max_select=2, mock=True)
        assert sum(p["keep"] for p in out["photos"]) == 2


class TestBuildDeck:
    def test_only_kept_become_slides(self, tmp_path):
        analysis = {
            "title": "測試相簿",
            "photos": [
                {"path": str(tmp_path / "a.jpg"), "keep": True, "caption": "海邊"},
                {"path": str(tmp_path / "b.jpg"), "keep": False, "caption": ""},
                {"path": str(tmp_path / "c.jpg"), "keep": True, "caption": "夕陽"},
            ],
        }
        _photo(tmp_path / "a.jpg"); _photo(tmp_path / "c.jpg")
        deck = pd.build_photo_deck(analysis, asset_base=tmp_path)
        slides = deck["sections"][0]["slides"]
        assert deck["deck_title"] == "測試相簿"
        assert deck["source_type"] == "photos"
        assert len(slides) == 2
        assert slides[0]["bg_image"] == "a.jpg" and slides[0]["narration"] == "海邊"
        assert slides[1]["bg_image"] == "c.jpg" and slides[1]["narration"] == "夕陽"
        assert all(s["reviewed"] is False for s in slides)

    def test_photos_to_deck_end_to_end_mock(self, tmp_path):
        paths = [_photo(tmp_path / f"p{i}.jpg") for i in range(3)]
        deck = pd.photos_to_deck(paths, asset_base=tmp_path, deck_title_hint="家庭日", mock=True)
        assert deck["deck_title"] == "家庭日"
        assert len(deck["sections"][0]["slides"]) == 3


class TestFeedsPptxExporter:
    def test_photo_deck_exports_pptx(self, tmp_path):
        pytest.importorskip("pptx", reason="需要 python-pptx")
        from core.slide_pptx import deck_to_pptx
        paths = [_photo(tmp_path / f"p{i}.jpg") for i in range(2)]
        deck = pd.photos_to_deck(paths, asset_base=tmp_path, deck_title_hint="旅遊", mock=True)
        out = deck_to_pptx(deck, tmp_path / "out.pptx", asset_base=tmp_path)
        assert out.exists()
        from pptx import Presentation
        prs = Presentation(str(out))
        assert len(prs.slides) == 2
        # 每張照片鋪成一張投影片 (至少一張圖)
        assert any(sh.shape_type == 13 for sh in prs.slides[0].shapes)
