"""core.pptx_augment — 在原始 .pptx 上就地補圖 (原文字保持可編輯) 測試。

不需 LibreOffice: augment_pptx 以 page_pngs 注入逐頁圖, 跳過渲染。需要 python-pptx
與 Pillow。AI 圖走 mock (PIL 佔位)。
"""
from __future__ import annotations

import pytest

pytest.importorskip("pptx", reason="需要 python-pptx")
pytest.importorskip("PIL", reason="需要 Pillow")

from core import pptx_augment as pa


def _make_pptx(path, n=2):
    """造一份有可編輯文字方塊的 pptx (16:9), 每頁左側放標題+內文。"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1.5))
        tb.text_frame.text = f"投影片 {i+1} 標題\n第 {i+1} 頁的內文說明"
    prs.save(str(path))
    return path


def _page_png(path, size=(960, 720)):
    """造一張「左半有內容、右半留白」的頁面圖 (利於偵測空白)。"""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", size, (245, 245, 245))
    ImageDraw.Draw(img).rectangle([20, 20, size[0] // 2 - 20, size[1] - 20], fill=(20, 20, 20))
    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path)
    return path


def _count(prs_path):
    from pptx import Presentation
    prs = Presentation(str(prs_path))
    res = []
    for slide in prs.slides:
        nt = sum(1 for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip())
        npic = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        res.append((nt, npic))
    return res


class TestExtractTexts:
    def test_extract(self, tmp_path):
        from pptx import Presentation
        src = _make_pptx(tmp_path / "d.pptx", n=2)
        texts = pa.extract_pptx_slide_texts(Presentation(str(src)))
        assert len(texts) == 2
        assert texts[0][0] == "投影片 1 標題"
        assert "內文" in texts[0][1]


class TestInsertImages:
    def test_inserts_and_preserves_text(self, tmp_path):
        src = _make_pptx(tmp_path / "d.pptx", n=2)
        ai = _page_png(tmp_path / "ai.png", size=(512, 512))
        out = tmp_path / "out.pptx"
        n = pa.insert_images_into_pptx(
            src, out,
            [(0, ai, (0.6, 0.1, 0.3, 0.4)), (1, ai, None)],
        )
        assert n == 2
        counts = _count(out)
        # 原文字方塊保留 (各 1), 各加 1 張圖
        assert counts == [(1, 1), (1, 1)]

    def test_skips_bad_index(self, tmp_path):
        src = _make_pptx(tmp_path / "d.pptx", n=1)
        ai = _page_png(tmp_path / "ai.png")
        out = tmp_path / "out.pptx"
        n = pa.insert_images_into_pptx(src, out, [(9, ai, None)])
        assert n == 0


class TestAugmentPptx:
    def test_augment_with_injected_pages_mock(self, tmp_path):
        src = _make_pptx(tmp_path / "deck.pptx", n=2)
        pages = [_page_png(tmp_path / "pages" / f"p{i:03d}.png") for i in (1, 2)]
        out = tmp_path / "aug.pptx"
        summary = pa.augment_pptx(
            src, out, work_dir=tmp_path / "work",
            only_missing=True, mock=True, page_pngs=pages,
        )
        assert summary["generated"] == 2
        assert summary["inserted"] == 2
        # 文字保留 + 每頁多一張圖
        assert _count(out) == [(1, 1), (1, 1)]

    def test_max_images_caps(self, tmp_path):
        src = _make_pptx(tmp_path / "deck.pptx", n=2)
        pages = [_page_png(tmp_path / "pages" / f"p{i:03d}.png") for i in (1, 2)]
        out = tmp_path / "aug.pptx"
        summary = pa.augment_pptx(
            src, out, work_dir=tmp_path / "work",
            mock=True, page_pngs=pages, max_images=1,
        )
        assert summary["inserted"] == 1


class TestPptxConverterBoundary:
    """CI 鎖 converter 選路；真正 Office round-trip 由 office_live gate 負責。"""

    def test_windows_without_soffice_uses_powerpoint(self, tmp_path, monkeypatch):
        src = tmp_path / "deck.pptx"
        src.write_bytes(b"PK")
        expected = tmp_path / "out" / "deck.pdf"
        calls = []

        monkeypatch.setattr(pa, "_IS_WINDOWS", True)
        monkeypatch.setattr(pa.shutil, "which", lambda _: None)

        def fake_powerpoint(source, output):
            calls.append((source, output))
            output.parent.mkdir(parents=True, exist_ok=True)
            output.write_bytes(b"%PDF-1.4")
            return output

        monkeypatch.setattr(pa, "_render_with_powerpoint", fake_powerpoint)
        result = pa.render_pptx_to_pdf(src, expected.parent)

        assert result == expected
        assert calls == [(src, expected)]

    def test_non_windows_without_converter_fails_actionably(self, tmp_path, monkeypatch):
        src = tmp_path / "deck.pptx"
        src.write_bytes(b"PK")
        monkeypatch.setattr(pa, "_IS_WINDOWS", False)
        monkeypatch.setattr(pa.shutil, "which", lambda _: None)

        with pytest.raises(RuntimeError, match="libreoffice-impress"):
            pa.render_pptx_to_pdf(src, tmp_path / "out")
