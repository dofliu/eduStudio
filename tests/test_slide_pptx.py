"""core.slide_pptx.deck_to_pptx + GET /jobs/{id}/pptx 測試。

涵蓋:
- 補過圖的頁 (image_generated) 放原頁+配圖兩張圖 + 旁白進備忘稿
- 沒補圖的頁放 bg_image
- slide 數 = deck slide 總數
- 路由: 有 deck → 回 .pptx; 無 deck → 404

需要 python-pptx 與 Pillow; 沒裝則 skip。
"""
from __future__ import annotations

import json

import pytest

pytest.importorskip("pptx", reason="需要 python-pptx")
pytest.importorskip("PIL", reason="需要 Pillow")

from core.slide_pptx import deck_to_pptx


def _png(path, size=(960, 720), color=(40, 40, 40)):
    from PIL import Image
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, color).save(path)


def _deck(base):
    _png(base / "p001.png")            # 補圖頁的原頁
    _png(base / "figures/ai_ch1_p001.png", size=(1024, 1024))  # AI 配圖
    _png(base / "figures/aug_ch1_p001.png", size=(1920, 1080))  # 合成頁
    _png(base / "p002.png")            # 沒補圖頁
    return {
        "deck_title": "測試簡報",
        "image_augmentation": {"layout": "side_by_side", "generated": 1},
        "sections": [{
            "id": "ch1", "title": "全部",
            "slides": [
                {"id": "ch1_p001", "title": "投影片 1", "narration": "第一頁旁白。",
                 "image_generated": True,
                 "source_bg_image": "p001.png",
                 "ai_image": "figures/ai_ch1_p001.png",
                 "bg_image": "figures/aug_ch1_p001.png"},
                {"id": "ch1_p002", "title": "投影片 2", "narration": "第二頁旁白。",
                 "bg_image": "p002.png"},
            ],
        }],
    }


class TestDeckToPptx:
    def test_exports_slides_and_notes(self, tmp_path):
        deck = _deck(tmp_path)
        out = deck_to_pptx(deck, tmp_path / "out.pptx", asset_base=tmp_path)
        assert out.exists() and out.stat().st_size > 0

        from pptx import Presentation
        prs = Presentation(str(out))
        assert len(prs.slides) == 2
        # 補圖頁有 2 張圖 (原頁+配圖); 旁白進備忘稿
        s1 = prs.slides[0]
        pics = [sh for sh in s1.shapes if sh.shape_type == 13]  # 13 = PICTURE
        assert len(pics) == 2
        assert "第一頁旁白" in s1.notes_slide.notes_text_frame.text
        # 沒補圖頁 1 張圖
        s2 = prs.slides[1]
        assert len([sh for sh in s2.shapes if sh.shape_type == 13]) == 1

    def test_overlay_layout(self, tmp_path):
        deck = _deck(tmp_path)
        deck["image_augmentation"]["layout"] = "overlay"
        out = deck_to_pptx(deck, tmp_path / "ov.pptx", asset_base=tmp_path)
        from pptx import Presentation
        prs = Presentation(str(out))
        # overlay 補圖頁仍是原頁 + 配圖兩張
        assert len([sh for sh in prs.slides[0].shapes if sh.shape_type == 13]) == 2


# ---------- route ----------

pytest.importorskip("fastapi.testclient", reason="需要 fastapi")

from fastapi.testclient import TestClient  # noqa: E402

import server.jobs as jobs_mod  # noqa: E402
from server.jobs import JobStore, get_default_store  # noqa: E402
from server.main import create_app  # noqa: E402
from server.schemas import CreateJobRequest, JobOptions, JobSource, SourceType  # noqa: E402


@pytest.fixture
def client(tmp_path, monkeypatch):
    jobs_root = tmp_path / "jobs"
    monkeypatch.setattr(jobs_mod, "JOBS_DIR", jobs_root)
    monkeypatch.setattr("core.config.PROJECT_ROOT", tmp_path)
    app = create_app()
    store = JobStore(root=jobs_root)
    app.dependency_overrides[get_default_store] = lambda: store
    with TestClient(app) as c:
        yield c, store, tmp_path


class TestPptxRoute:
    def test_download_pptx(self, client):
        c, store, base = client
        rec = store.create(CreateJobRequest(
            source_type=SourceType.SLIDES_PDF,
            source=JobSource(path=str(base / "x.pdf")),
            options=JobOptions(),
        ))
        # 圖片相對 PROJECT_ROOT(=tmp_path)
        deck = _deck(base)
        store.deck_path(rec.id).write_text(json.dumps(deck), encoding="utf-8")

        r = c.get(f"/jobs/{rec.id}/pptx")
        assert r.status_code == 200, r.text
        assert "presentationml" in r.headers["content-type"]
        assert r.content[:2] == b"PK"  # .pptx = zip

    def test_no_deck_404(self, client):
        c, store, base = client
        rec = store.create(CreateJobRequest(
            source_type=SourceType.SLIDES_PDF,
            source=JobSource(path=str(base / "x.pdf")),
            options=JobOptions(),
        ))
        r = c.get(f"/jobs/{rec.id}/pptx")
        assert r.status_code == 404
