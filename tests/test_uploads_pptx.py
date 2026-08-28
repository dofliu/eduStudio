"""POST /upload/pptx — PPTX 原檔就地補圖路由測試。

驗證: 非 .pptx → 400; 空檔 → 400; happy path (建 job + 背景補圖, mock) — happy
path 需 LibreOffice (pptx→pdf), 無 soffice 則 skip。
"""
from __future__ import annotations

import io
import importlib.util
import os
import shutil

import pytest

pytest.importorskip("fastapi.testclient", reason="需要 fastapi")
pytest.importorskip("multipart", reason="需要 python-multipart")
pytest.importorskip("pptx", reason="需要 python-pptx")

from fastapi.testclient import TestClient

import server.jobs as jobs_mod
from server.jobs import JobStore, get_default_store
from server.main import create_app

_HAS_SOFFICE = bool(shutil.which("soffice") or shutil.which("libreoffice"))
_HAS_POWERPOINT = False
if os.name == "nt" and importlib.util.find_spec("win32com") is not None:
    try:
        import winreg
        with winreg.OpenKey(winreg.HKEY_CLASSES_ROOT, r"PowerPoint.Application\CurVer"):
            _HAS_POWERPOINT = True
    except OSError:
        pass
_HAS_PPT_CONVERTER = _HAS_SOFFICE or _HAS_POWERPOINT


def _pptx_bytes(n=2):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1)).text_frame.text = f"頁 {i+1}"
    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()


@pytest.fixture
def client(tmp_path, monkeypatch):
    jobs_root = tmp_path / "jobs"
    monkeypatch.setattr(jobs_mod, "JOBS_DIR", jobs_root)
    app = create_app()
    store = JobStore(root=jobs_root)
    app.dependency_overrides[get_default_store] = lambda: store
    with TestClient(app) as c:
        yield c, store


class TestValidation:
    def test_non_pptx_rejected(self, client):
        c, _ = client
        r = c.post("/upload/pptx", files={"file": ("x.pdf", b"%PDF", "application/pdf")})
        assert r.status_code == 400

    def test_empty_rejected(self, client):
        c, _ = client
        r = c.post("/upload/pptx", files={"file": ("x.pptx", b"", "application/vnd.openxmlformats-officedocument.presentationml.presentation")})
        assert r.status_code == 400


class TestToVideo:
    """POST /jobs/{id}/to-video — 補圖 pptx → slides_pdf 影片 job。"""

    def _pptx_job_with_artifact(self, store):
        from server.schemas import CreateJobRequest, JobOptions, JobSource, SourceType
        rec = store.create(CreateJobRequest(
            source_type=SourceType.PPTX, source=JobSource(), options=JobOptions(mock=True)))
        art = store.artifacts_dir(rec.id) / "deck_augmented.pptx"
        art.parent.mkdir(parents=True, exist_ok=True)
        art.write_bytes(b"PKfake")
        store.refresh_artifacts(rec.id)
        return rec

    def test_creates_slides_job(self, client, monkeypatch):
        c, store = client
        rec = self._pptx_job_with_artifact(store)

        import core.pptx_augment as pa
        import server.routes.jobs as jr

        def fake_render(src, outdir, **k):
            from pathlib import Path
            p = Path(outdir); p.mkdir(parents=True, exist_ok=True)
            pdf = p / "deck.pdf"; pdf.write_bytes(b"%PDF-1.4"); return pdf
        monkeypatch.setattr(pa, "render_pptx_to_pdf", fake_render)
        monkeypatch.setattr(jr, "schedule_job", lambda store, jid: None)

        r = c.post(f"/jobs/{rec.id}/to-video")
        assert r.status_code == 201, r.text
        new_id = r.json()["job_id"]
        assert store.get(new_id).source_type.value == "slides_pdf"

    def test_non_pptx_rejected(self, client):
        c, store = client
        from server.schemas import CreateJobRequest, JobOptions, JobSource, SourceType
        rec = store.create(CreateJobRequest(
            source_type=SourceType.SLIDES_PDF, source=JobSource(path="x.pdf"),
            options=JobOptions()))
        assert c.post(f"/jobs/{rec.id}/to-video").status_code == 400

    def test_no_artifact_rejected(self, client):
        c, store = client
        from server.schemas import CreateJobRequest, JobOptions, JobSource, SourceType
        rec = store.create(CreateJobRequest(
            source_type=SourceType.PPTX, source=JobSource(), options=JobOptions()))
        assert c.post(f"/jobs/{rec.id}/to-video").status_code == 400


@pytest.mark.skipif(not _HAS_PPT_CONVERTER, reason="需要 LibreOffice 或 PowerPoint 做 pptx→pdf")
class TestHappyPath:
    def test_creates_job_and_augments(self, client):
        c, store = client
        r = c.post(
            "/upload/pptx",
            files={"file": ("deck.pptx", _pptx_bytes(2), "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            data={"only_missing": "true", "options_json": '{"mock": true}'},
        )
        assert r.status_code == 201, r.text
        job_id = r.json()["job_id"]
        assert store.get(job_id).source_type.value == "pptx"
        # 背景 task 在 TestClient 結束時應已跑完 (mock 很快); 輪詢 deck 不適用, 檢查 state
        import time
        for _ in range(150):  # LibreOffice 轉檔在多測試同跑時可能較慢, 給足 30s
            rec = store.get(job_id)
            if rec.state.value in ("done", "failed"):
                break
            time.sleep(0.2)
        rec = store.get(job_id)
        assert rec.state.value == "done", rec.error
        names = [a.name for a in rec.artifacts]
        assert any(n.endswith("_augmented.pptx") for n in names), names
